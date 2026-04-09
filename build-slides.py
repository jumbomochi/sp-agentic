"""
Build the full-day workshop slide deck.

Uses existing Google course PPTX as template source for backgrounds,
layouts, and theme. Creates slides for all sessions including segue
slides for sections covered by the official Google materials.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy
import io
import os

# ── Colors (Google palette) ─────────────────────────────────────
BLUE   = RGBColor(0x42, 0x85, 0xF4)
RED    = RGBColor(0xEA, 0x43, 0x35)
YELLOW = RGBColor(0xFB, 0xBC, 0x05)
GREEN  = RGBColor(0x34, 0xA8, 0x53)
DARK   = RGBColor(0x20, 0x21, 0x24)
MID    = RGBColor(0x3C, 0x40, 0x43)
SOFT   = RGBColor(0x5F, 0x63, 0x68)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
ACCENT_BG = RGBColor(0xE8, 0xF0, 0xFE)

# ── Dimensions (matching Google Slides 20"x11.25") ──────────────
SLIDE_W = Emu(18288000)
SLIDE_H = Emu(10287000)

# Margins
LEFT_M = Emu(623400)      # ~0.68"
TOP_M  = Emu(890050)       # ~0.97"
CONTENT_W = Emu(17041200)  # ~18.6"

# ── Font names ──────────────────────────────────────────────────
FONT_TITLE = "Google Sans"
FONT_BODY  = "Google Sans"
FONT_MONO  = "Roboto Mono"

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
GW_DIR = os.path.join(BASE_DIR, "Google Workspace with Gemini")
AD_DIR = os.path.join(BASE_DIR, "Application Development with LLMs on Google Cloud")


def load_template():
    """Load an existing PPTX as template to get the slide master + theme."""
    template_path = os.path.join(
        GW_DIR,
        "Google Workspace with Gemini - Module 0 - Course Introduction.pptx"
    )
    prs = Presentation(template_path)
    # Clear all existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    return prs


def get_layout(prs, name):
    """Get a slide layout by name."""
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    # Fallback
    return prs.slide_layouts[0]


def extract_bg_image(source_pptx_path):
    """Extract the full-bleed background image from slide 1 of a PPTX."""
    src = Presentation(source_pptx_path)
    slide = list(src.slides)[0]
    for shape in slide.shapes:
        if shape.shape_type == 13 and shape.width > 10000000:
            return shape.image.blob, shape.image.content_type
    return None, None


def add_bg_image(slide, img_blob, content_type="image/png"):
    """Add a full-bleed background image to a slide."""
    from pptx.util import Emu
    slide.shapes.add_picture(
        io.BytesIO(img_blob),
        Emu(0), Emu(0),
        SLIDE_W, SLIDE_H
    )
    # Move it to back
    sp = slide.shapes[-1]._element
    slide.shapes._spTree.remove(sp)
    slide.shapes._spTree.insert(2, sp)


def add_text_box(slide, left, top, width, height, text,
                 font_name=FONT_BODY, font_size=Pt(28), font_color=DARK,
                 bold=False, alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 line_spacing=2.0):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    p.space_after = Pt(6)
    p.line_spacing = line_spacing

    return txBox


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=Pt(24), font_color=MID, spacing=2.0):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.name = FONT_BODY
        p.font.size = font_size
        p.font.color.rgb = font_color
        p.line_spacing = spacing
        p.space_after = Pt(8)
        p.level = 0
        # Add bullet character
        p.text = "•  " + item

    return txBox


def add_number_badge(slide, number, color=BLUE, left=None, top=None):
    """Add a numbered badge like the Google slides use."""
    if left is None:
        left = Emu(16300000)
    if top is None:
        top = Emu(500000)

    # Circle background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left, top,
        Emu(1200000), Emu(1200000)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    # Number text
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = f"{number:02d}"
    p.font.name = FONT_MONO
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def add_colored_bar(slide, color, left=Emu(0), top=Emu(0),
                    width=Emu(18288000), height=Emu(120000)):
    """Add a thin colored bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ═══════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════

def build_title_slide(prs, bg_blob):
    """Slide 1: Course title slide."""
    layout = get_layout(prs, "BLANK")
    slide = prs.slides.add_slide(layout)

    if bg_blob:
        add_bg_image(slide, bg_blob)

    # Title
    add_text_box(
        slide, Emu(1000000), Emu(2500000), Emu(16000000), Emu(3000000),
        "AI-Powered Productivity with\nGoogle Gemini & LLM Development",
        font_name=FONT_TITLE, font_size=Pt(54), font_color=DARK,
        bold=False, alignment=PP_ALIGN.CENTER
    )

    # Subtitle
    add_text_box(
        slide, Emu(1000000), Emu(5800000), Emu(16000000), Emu(1500000),
        "One-Day Workshop",
        font_name=FONT_BODY, font_size=Pt(32), font_color=SOFT,
        alignment=PP_ALIGN.CENTER
    )

    # Audience
    add_text_box(
        slide, Emu(1000000), Emu(7200000), Emu(16000000), Emu(1200000),
        "Singapore Polytechnic  •  School of Mathematical Sciences and Analytics",
        font_name=FONT_BODY, font_size=Pt(24), font_color=SOFT,
        alignment=PP_ALIGN.CENTER
    )

    return slide


def build_section_header(prs, bg_blob, session_num, title, subtitle="", color=BLUE):
    """Section header slide with number badge."""
    layout = get_layout(prs, "BLANK")
    slide = prs.slides.add_slide(layout)

    if bg_blob:
        add_bg_image(slide, bg_blob)

    # Number badge (top-right)
    add_number_badge(slide, session_num, color=color,
                     left=Emu(16300000), top=Emu(500000))

    # Title (right-aligned area, like Google template)
    add_text_box(
        slide, Emu(9200000), Emu(3200000), Emu(8000000), Emu(2500000),
        title,
        font_name=FONT_TITLE, font_size=Pt(48), font_color=DARK,
        bold=False
    )

    if subtitle:
        add_text_box(
            slide, Emu(9200000), Emu(5500000), Emu(8000000), Emu(1500000),
            subtitle,
            font_name=FONT_BODY, font_size=Pt(24), font_color=SOFT
        )

    return slide


def build_content_slide(prs, title, bullets, footer_text=""):
    """Standard content slide with title and bullets."""
    layout = get_layout(prs, "BLANK")
    slide = prs.slides.add_slide(layout)

    # Light background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Colored top bar
    add_colored_bar(slide, BLUE, height=Emu(80000))

    # Title
    add_text_box(
        slide, LEFT_M, Emu(500000), CONTENT_W, Emu(1200000),
        title,
        font_name=FONT_TITLE, font_size=Pt(40), font_color=DARK
    )

    # Divider line
    add_colored_bar(
        slide, RGBColor(0xDA, 0xDC, 0xE0),
        left=LEFT_M, top=Emu(1600000),
        width=CONTENT_W, height=Emu(20000)
    )

    # Bullets
    add_bullet_list(
        slide, LEFT_M, Emu(1900000), Emu(16000000), Emu(7000000),
        bullets, font_size=Pt(26), font_color=MID
    )

    # Footer
    if footer_text:
        add_text_box(
            slide, LEFT_M, Emu(9400000), CONTENT_W, Emu(600000),
            footer_text,
            font_name=FONT_BODY, font_size=Pt(16), font_color=SOFT,
            alignment=PP_ALIGN.LEFT
        )

    return slide


def build_two_column_slide(prs, title, left_items, right_items,
                           left_header="", right_header=""):
    """Two-column comparison slide."""
    layout = get_layout(prs, "BLANK")
    slide = prs.slides.add_slide(layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    add_colored_bar(slide, BLUE, height=Emu(80000))

    # Title
    add_text_box(
        slide, LEFT_M, Emu(500000), CONTENT_W, Emu(1200000),
        title,
        font_name=FONT_TITLE, font_size=Pt(40), font_color=DARK
    )

    add_colored_bar(
        slide, RGBColor(0xDA, 0xDC, 0xE0),
        left=LEFT_M, top=Emu(1600000),
        width=CONTENT_W, height=Emu(20000)
    )

    col_width = Emu(7800000)
    gap = Emu(800000)

    # Left column header
    if left_header:
        add_text_box(
            slide, LEFT_M, Emu(1900000), col_width, Emu(600000),
            left_header,
            font_name=FONT_TITLE, font_size=Pt(28), font_color=BLUE, bold=True
        )

    add_bullet_list(
        slide, LEFT_M, Emu(2500000), col_width, Emu(6500000),
        left_items, font_size=Pt(22), font_color=MID
    )

    # Right column header
    right_left = Emu(623400) + col_width + gap
    if right_header:
        add_text_box(
            slide, right_left, Emu(1900000), col_width, Emu(600000),
            right_header,
            font_name=FONT_TITLE, font_size=Pt(28), font_color=GREEN, bold=True
        )

    add_bullet_list(
        slide, right_left, Emu(2500000), col_width, Emu(6500000),
        right_items, font_size=Pt(22), font_color=MID
    )

    # Vertical divider
    add_colored_bar(
        slide, RGBColor(0xDA, 0xDC, 0xE0),
        left=Emu(623400) + col_width + Emu(350000),
        top=Emu(1900000),
        width=Emu(20000), height=Emu(7000000)
    )

    return slide


def build_segue_slide(prs, bg_blob, text, subtext="", color=BLUE):
    """Transition slide pointing to official Google materials."""
    layout = get_layout(prs, "BLANK")
    slide = prs.slides.add_slide(layout)

    if bg_blob:
        add_bg_image(slide, bg_blob)

    add_text_box(
        slide, Emu(2000000), Emu(3500000), Emu(14000000), Emu(2000000),
        text,
        font_name=FONT_TITLE, font_size=Pt(44), font_color=DARK,
        alignment=PP_ALIGN.CENTER
    )

    if subtext:
        add_text_box(
            slide, Emu(2000000), Emu(5500000), Emu(14000000), Emu(1500000),
            subtext,
            font_name=FONT_BODY, font_size=Pt(24), font_color=SOFT,
            alignment=PP_ALIGN.CENTER
        )

    return slide


def build_break_slide(prs, bg_blob, text="Break", duration="15 min"):
    """Break / lunch slide."""
    layout = get_layout(prs, "BLANK")
    slide = prs.slides.add_slide(layout)

    if bg_blob:
        add_bg_image(slide, bg_blob)

    add_text_box(
        slide, Emu(2000000), Emu(3500000), Emu(14000000), Emu(2000000),
        text,
        font_name=FONT_TITLE, font_size=Pt(54), font_color=DARK,
        alignment=PP_ALIGN.CENTER
    )

    add_text_box(
        slide, Emu(2000000), Emu(5500000), Emu(14000000), Emu(1000000),
        duration,
        font_name=FONT_BODY, font_size=Pt(28), font_color=SOFT,
        alignment=PP_ALIGN.CENTER
    )

    return slide


def build_exercise_slide(prs, title, instructions, color=BLUE):
    """Hands-on exercise slide."""
    layout = get_layout(prs, "BLANK")
    slide = prs.slides.add_slide(layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    add_colored_bar(slide, color, height=Emu(80000))

    # Exercise icon/label
    add_text_box(
        slide, LEFT_M, Emu(500000), Emu(3000000), Emu(600000),
        "Hands-on Exercise",
        font_name=FONT_TITLE, font_size=Pt(22), font_color=color, bold=True
    )

    # Title
    add_text_box(
        slide, LEFT_M, Emu(1100000), CONTENT_W, Emu(1000000),
        title,
        font_name=FONT_TITLE, font_size=Pt(38), font_color=DARK
    )

    add_colored_bar(
        slide, RGBColor(0xDA, 0xDC, 0xE0),
        left=LEFT_M, top=Emu(2000000),
        width=CONTENT_W, height=Emu(20000)
    )

    add_bullet_list(
        slide, LEFT_M, Emu(2300000), Emu(16000000), Emu(7000000),
        instructions, font_size=Pt(24), font_color=MID
    )

    return slide


def build_table_slide(prs, title, headers, rows):
    """Slide with a table."""
    layout = get_layout(prs, "BLANK")
    slide = prs.slides.add_slide(layout)

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    add_colored_bar(slide, BLUE, height=Emu(80000))

    add_text_box(
        slide, LEFT_M, Emu(500000), CONTENT_W, Emu(1000000),
        title,
        font_name=FONT_TITLE, font_size=Pt(40), font_color=DARK
    )

    # Build table
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header

    table_left = LEFT_M
    table_top = Emu(1800000)
    table_width = Emu(16000000)
    table_height = Emu(num_rows * 800000)

    table_shape = slide.shapes.add_table(
        num_rows, num_cols, table_left, table_top, table_width, table_height
    )
    table = table_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            para.font.name = FONT_TITLE
            para.font.size = Pt(20)
            para.font.color.rgb = WHITE
            para.font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                para.font.name = FONT_BODY
                para.font.size = Pt(18)
                para.font.color.rgb = MID
            # Alternate row colors
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return slide


# ═══════════════════════════════════════════════════════════════
# MAIN: Build the full deck
# ═══════════════════════════════════════════════════════════════

def main():
    print("Loading template...")
    prs = load_template()

    # Extract background images from existing presentations
    print("Extracting background images...")
    gw_title_bg, _ = extract_bg_image(os.path.join(
        GW_DIR, "Google Workspace with Gemini - Module 0 - Course Introduction.pptx"))
    gw_section_bg, _ = extract_bg_image(os.path.join(
        GW_DIR, "Google Workspace with Gemini - Module 2 - Gemini app.pptx"))
    ad_section_bg, _ = extract_bg_image(os.path.join(
        AD_DIR, "M2-Vertex-AI-Studio.pptx"))

    # ── SLIDE 1: Title ─────────────────────────────────────────
    print("Building slides...")
    build_title_slide(prs, gw_title_bg)

    # ── SESSION 1: Welcome & Gemini Ecosystem ──────────────────
    build_section_header(prs, gw_section_bg, 1,
        "Welcome &\nThe Gemini Ecosystem",
        "9:00 AM – 9:30 AM", color=BLUE)

    build_content_slide(prs,
        "Today's Journey",
        [
            "From using AI as a daily productivity tool...",
            "...to developing LLM-powered applications on Google Cloud",
            "Hands-on exercises at every stage — beginner and intermediate tracks",
            "No prior coding experience required for most sessions",
            "By end of day: Gems, Apps Script, Workspace Flows, Vertex AI, Gemini API",
        ]
    )

    build_table_slide(prs,
        "The Google Gemini Ecosystem",
        ["Product", "What it is", "Best for"],
        [
            ["Gemini app", "Standalone AI assistant", "Brainstorming, deep research, image/video generation"],
            ["Workspace with Gemini", "AI in Gmail, Docs, Sheets, Slides, Meet, Drive", "Enhancing daily productivity in apps you already use"],
            ["NotebookLM", "AI grounded in your documents", "Synthesising insights without hallucination"],
            ["Google Vids", "AI-powered video creation", "Training videos, presentations, announcements"],
            ["Vertex AI", "Google Cloud's AI/ML platform", "Building LLM apps, prompt testing, Gemini API"],
        ]
    )

    build_content_slide(prs,
        "From Prompts to Applications — The AI Maturity Spectrum",
        [
            "Level 1 — Prompting: Ask Gemini to draft an email about the CA test",
            "Level 2 — Customising: Create a Gem with persona and instructions for reuse",
            "Level 3 — Code Automation: Apps Script sends personalised grade emails",
            "Level 4 — No-Code Pipelines: Workspace Flows compose Ask Gemini with Gmail, Chat, Sheets",
            "Level 5 — Developing Apps: Gemini API via Python for custom applications",
        ],
        footer_text="Today we touch every level on this spectrum."
    )

    build_content_slide(prs,
        "Using Gemini Responsibly in Education",
        [
            "Hallucination: LLMs can present incorrect information confidently — always verify",
            "Grounding: Constraining AI to known documents reduces hallucination (NotebookLM, Flows, Vertex AI)",
            "Academic integrity: AI is a tool for educators, not a shortcut for students",
            "Data privacy: Do NOT paste student personal data into the public Gemini app",
            "Google Workspace with Gemini has enterprise data protections — but check your institution's policy",
        ]
    )

    # ── SESSION 2: Segue to Google Workspace — Gemini App & Gems
    build_section_header(prs, gw_section_bg, 2,
        "The Gemini App &\nCustom Gems",
        "9:30 AM – 10:15 AM", color=BLUE)

    build_segue_slide(prs, gw_title_bg,
        "Switching to Google Workspace with Gemini slides",
        "Module 2 — Gemini app: Deep Research, Canvas, Gems, Knowledge files"
    )

    # ── BREAK ──────────────────────────────────────────────────
    build_break_slide(prs, gw_title_bg, "Break", "10:15 AM – 10:30 AM")

    # ── SESSION 3: Segue to Sheets ────────────────────────────
    build_section_header(prs, gw_section_bg, 3,
        "Gemini in Google Sheets\nAnalytics Focus",
        "10:30 AM – 11:15 AM", color=GREEN)

    build_segue_slide(prs, gw_title_bg,
        "Switching to Google Workspace with Gemini slides",
        "Module 6 — Gemini in Google Sheets: Tables, formulas, charts, =AI() function"
    )

    # ── SESSION 4: Segue to Docs & Gmail ──────────────────────
    build_section_header(prs, gw_section_bg, 4,
        "Gemini in Docs & Gmail\nTeaching Materials & Admin",
        "11:15 AM – 11:45 AM", color=GREEN)

    build_segue_slide(prs, gw_title_bg,
        "Switching to Google Workspace with Gemini slides",
        "Modules 4 & 3 — Gemini in Google Docs + Gemini in Gmail"
    )

    # ── SESSION 5: NotebookLM Demo ────────────────────────────
    build_section_header(prs, gw_section_bg, 5,
        "NotebookLM for\nResearch & Teaching",
        "11:45 AM – 12:00 PM  •  Demo only", color=GREEN)

    build_content_slide(prs,
        "What is NotebookLM?",
        [
            "AI research assistant grounded entirely in your uploaded sources",
            "Unlike Gemini: answers only from your documents — not from training data",
            "Supports Google Docs, Slides, PDFs, websites, YouTube videos, audio files",
            "Up to 50 sources per notebook, 500,000 words each",
            "Inline citations link every claim to a specific source passage",
        ]
    )

    build_two_column_slide(prs,
        "NotebookLM vs. Gems",
        [
            "Custom prompt template",
            "Upload static knowledge files",
            "Text output only",
            "No citations",
            "Best for: repetitive tasks with consistent format",
        ],
        [
            "Grounded in your specific documents",
            "Inline citations for every claim",
            "Audio overviews (podcast-style)",
            "Cross-source synthesis",
            "Best for: research, teaching prep, study hubs",
        ],
        left_header="Gems",
        right_header="NotebookLM"
    )

    build_content_slide(prs,
        "Live Demo: Statistics I Teaching Hub",
        [
            "Create a notebook with module FAQ + textbook chapter",
            "Query 1: \"What is the grading breakdown?\" → single-source with citations",
            "Query 2: \"What concepts to focus on for the exam?\" → cross-source synthesis",
            "Query 3: \"Create a 5-question quiz on probability\" → grounded content generation",
            "Audio overview: AI-generated podcast summarising your materials",
        ],
        footer_text="Try it yourself: notebooklm.google.com"
    )

    # ── LUNCH ──────────────────────────────────────────────────
    build_break_slide(prs, gw_title_bg, "Lunch", "12:00 PM – 1:00 PM")

    # ── SESSION 6: Apps Script ────────────────────────────────
    build_section_header(prs, gw_section_bg, 6,
        "Apps Script Generation\nwith Gemini",
        "1:00 PM – 1:45 PM", color=YELLOW)

    build_content_slide(prs,
        "What is Google Apps Script?",
        [
            "JavaScript-based automation platform built into Google Workspace",
            "Already available in every Sheet, Doc, and Form — no installation needed",
            "Automates tasks across Workspace apps: Sheets ↔ Gmail ↔ Calendar ↔ Forms",
            "Runs on Google's servers — no local setup or hosting required",
            "You don't need to know JavaScript — Gemini generates the code for you",
        ]
    )

    build_content_slide(prs,
        "The Prompting Pattern for Apps Script",
        [
            "Act as a Google Apps Script developer.",
            "Read data from a Google Sheet named \"[sheet name]\"",
            "Describe the automation in plain language",
            "Include comments explaining each section",
            "Handle errors gracefully",
            "Add a custom menu item to run the script from the spreadsheet",
        ],
        footer_text="Be specific about sheet names, column layouts, and desired behaviour."
    )

    build_content_slide(prs,
        "Demo 1: Auto-Email Student Grades",
        [
            "Reads student data from a Google Sheet (ID, Name, Email, Scores)",
            "Sends personalised emails with grades and performance-based messages",
            "Encouraging for high performers, supportive for those struggling",
            "Writes \"Sent\" status back to the spreadsheet",
            "Custom menu: Grade Tools → Send Grade Emails",
        ],
        footer_text="Live demo: prompting Gemini → reviewing code → pasting into Script Editor → running"
    )

    build_content_slide(prs,
        "Demo 2: Automatic Form Response Processor",
        [
            "Triggers automatically when a student submits a Google Form",
            "Logs responses to a colour-coded \"Feedback Summary\" sheet",
            "Sends an acknowledgement email to the student",
            "Alerts the lecturer immediately for critical feedback (rating ≤ 2)",
            "No manual intervention — runs entirely on triggers",
        ],
        footer_text="Live demo: prompting Gemini → setting up trigger → submitting test form → observing automation"
    )

    build_exercise_slide(prs,
        "Choose Your Apps Script Scenario (15 min)",
        [
            "Beginner: Generate a script that sends personalised emails to each student in a Sheets list",
            "Intermediate: Generate a form-triggered script that logs responses and sends a summary email",
            "Steps: (1) Prepare data in a Sheet → (2) Prompt Gemini → (3) Paste into Script Editor → (4) Run and test",
            "Stuck? Use the starter prompts in the exercise handout",
            "Finished early? Add a time-based trigger or conditional formatting via script",
        ],
        color=YELLOW
    )

    # ── SESSION 7: Workspace Studio Flows ────────────────────
    build_section_header(prs, gw_section_bg, 7,
        "Workspace\nStudio Flows",
        "1:45 PM – 2:30 PM", color=YELLOW)

    build_content_slide(prs,
        "What are Workspace Studio Flows?",
        [
            "No-code, event-driven automations built in a visual pipeline",
            "Trigger → AI step → Actions — compose Gemini reasoning with Workspace apps",
            "13 triggers across Gmail, Chat, Sheets, Drive, Forms, Meet, Calendar, schedules",
            "20+ actions: Ask Gemini, Gmail, Chat, Sheets, Docs, Tasks, control flow",
            "Like Power Automate with Copilot steps — native to Google Workspace",
        ]
    )

    build_two_column_slide(prs,
        "Apps Script vs. Workspace Studio Flows",
        [
            "Code (JavaScript)",
            "Gemini generates the code for you",
            "Complex logic, full control",
            "AI calls require explicit API usage",
            "Example: Grade emailer script",
        ],
        [
            "No code — visual drag-and-configure",
            "Ask Gemini is a first-class step",
            "Event-driven pipelines across apps",
            "AI reasoning composed inline with actions",
            "Example: Intelligent Inbox Triage flow",
        ],
        left_header="Apps Script (Session 6)",
        right_header="Workspace Studio Flows (Session 7)"
    )

    build_content_slide(prs,
        "Live Build: Intelligent Inbox Triage",
        [
            "Trigger: When I get an email (filtered by [DEMO] subject)",
            "Step 1: Ask Gemini — classify + extract (custom prompt)",
            "Step 2: Check if — only continue for urgent cases",
            "Step 3: Notify me in Chat — formatted summary with extracted fields",
            "Step 4: Add labels — apply Urgent-Student to the email",
        ],
        footer_text="5 nodes. 1 prompt. 1 filter. No code — just a visual pipeline with Gemini in the middle."
    )

    build_content_slide(prs,
        "Writing the Ask Gemini Prompt",
        [
            "Define the role: \"You are an email triage assistant for a SP lecturer\"",
            "Give a closed set of categories — explicit choices beat open-ended ones",
            "Demand a rigid output format — every field on its own line",
            "The next step reads your output programmatically — format drift breaks the flow",
            "Iterate on the prompt, not on code — save, resend, watch it work",
        ],
        footer_text="Same prompt-engineering discipline as Gems (Session 2) and =AI() cells (Session 3)."
    )

    build_exercise_slide(prs,
        "Build Your Own Flow (20 min)",
        [
            "Beginner: Morning Inbox Digest (On a schedule → Recap unread emails → Notify in Chat)",
            "Intermediate: Email-to-Task Extractor (Ask Gemini extracts deadlines, creates Google Tasks)",
            "Advanced: Student Query Router (multi-branch classification and Chat routing)",
            "Fallback: Bring your own use case — describe it in the Discover bar",
            "Test-trigger your flow, iterate on the prompt, share if you like it",
        ],
        color=YELLOW
    )

    # ── BREAK ──────────────────────────────────────────────────
    build_break_slide(prs, gw_title_bg, "Break", "2:45 PM – 3:00 PM")

    # ── SESSION 8: Vertex AI & Prompt Engineering ─────────────
    build_section_header(prs, ad_section_bg, 8,
        "Vertex AI Studio &\nPrompt Engineering",
        "3:00 PM – 3:20 PM", color=RED)

    build_content_slide(prs,
        "What is Vertex AI?",
        [
            "Google Cloud's unified AI/ML platform",
            "Model Garden: catalogue of pre-trained models (Gemini, Llama, Mistral, and more)",
            "Vertex AI Studio: low-code prompt playground for testing and tuning",
            "API access: call Gemini models programmatically from Python, Node.js, Go, etc.",
            "The bridge from 'prompt in a chat box' to 'prompt in your application'",
        ]
    )

    build_content_slide(prs,
        "Prompt Engineering Techniques",
        [
            "Zero-shot: Ask the question directly — no examples provided",
            "One-shot / Few-shot: Provide 1 or more examples to guide the output format",
            "Chain-of-thought: \"Let's think step by step\" — improves reasoning accuracy",
            "System prompts: Set the model's persona, constraints, and output format",
            "Grounding: Provide documents or enable search for factual accuracy",
        ],
        footer_text="You've been using these all day — now you'll see them in Vertex AI Studio and the API."
    )

    build_content_slide(prs,
        "From Playground to Code",
        [
            "Vertex AI Studio: test prompts visually → tune parameters → see results",
            "\"Get Code\" button: exports your working prompt as Python, Node.js, or cURL",
            "The same system instructions and grounding concepts from Workspace apply here",
            "Lab 1 (next): explore the Vertex AI Studio UI — no coding required",
            "Lab 2 (after): use the Gemini API from Python notebooks",
        ]
    )

    # ── SESSION 9: Lab 1 — Vertex AI Studio UI ────────────────
    build_section_header(prs, ad_section_bg, 9,
        "Lab 1: Getting Started with\nVertex AI Studio UI",
        "3:20 PM – 3:50 PM  •  Google Cloud Skills Boost", color=RED)

    build_segue_slide(prs, ad_section_bg,
        "Switching to Google Cloud Skills Boost",
        "Lab: Getting Started with the Vertex AI Studio User Interface\nNo coding required — explore the prompt playground, test models, export code"
    )

    # ── SESSION 10: Lab 2 — Gemini API ────────────────────────
    build_section_header(prs, ad_section_bg, 10,
        "Lab 2: Gen AI Library +\nVertex AI Gemini API",
        "3:50 PM – 4:35 PM  •  Google Cloud Skills Boost", color=RED)

    build_content_slide(prs,
        "Lab 2 — Part 1: API, Chat & Grounding",
        [
            "Install google-genai and connect with genai.Client(vertexai=True, ...)",
            "Build a chat() helper wrapping generate_content — observe knowledge cutoff limits",
            "Manage chat_history as a list to enable multi-turn memory",
            "Add system messages to ground the model's persona and behaviour",
            "Upload a text file to Cloud Storage for document grounding",
            "Enable the GoogleSearch tool for real-time web information",
            "Explore the Google Maps tool for location-aware queries",
        ],
        footer_text="Code is pre-written — run cells and observe. Modify if you're comfortable."
    )

    build_content_slide(prs,
        "Lab 2 — Part 2: Embeddings & Similarity",
        [
            "Generate embeddings using the gemini-embedding-001 model",
            "Compare embeddings with L1 (Manhattan) and L2 (Euclidean) distance",
            "Use Cosine similarity and Dot product for semantic comparison",
            "Test: how similar are \"supercalifragilistic\" and its misspelling?",
            "Compute similarity between a user query and a document chunk",
            "This is the foundation of RAG — retrieve relevant chunks before prompting",
        ],
        footer_text="Embeddings + vector search = how production LLM apps ground answers in private data."
    )

    build_segue_slide(prs, ad_section_bg,
        "Switching to Google Cloud Skills Boost",
        "Lab: Getting Started with Gen AI Library + Vertex AI Gemini API\nintro_to_gemini-v1.0.0.ipynb — follow along step by step"
    )

    # ── SESSION 11: Wrap-Up ───────────────────────────────────
    build_section_header(prs, gw_section_bg, 11,
        "Wrap-Up &\nAction Planning",
        "4:35 PM – 4:50 PM", color=BLUE)

    build_table_slide(prs,
        "The Spectrum You Experienced Today",
        ["Level", "What you did", "Session"],
        [
            ["Using AI", "Prompted Gemini in Sheets, Docs, Gmail", "2–4"],
            ["Customising AI", "Created Gems with persona and instructions", "2"],
            ["Researching with AI", "Saw NotebookLM grounded in your documents", "5"],
            ["Automating with AI", "Generated Apps Script automations", "6"],
            ["Building with AI", "Built a no-code Workspace Flow", "7"],
            ["Developing with AI", "Used Vertex AI Studio and the Gemini API", "9–10"],
        ]
    )

    build_content_slide(prs,
        "The Thread: Grounding",
        [
            "Gems: grounded with knowledge files and instructions",
            "NotebookLM: entirely grounded in your uploaded sources",
            "Workspace Flows: the Ask Gemini step is grounded in the trigger event data",
            "Apps Script: grounded in your spreadsheet data",
            "Vertex AI Studio: system prompts and document grounding",
            "Gemini API: same grounding concepts applied in code",
        ],
        footer_text="Every time we constrained the model, the output got better and more reliable."
    )

    build_content_slide(prs,
        "Quick Wins for Monday",
        [
            "1. Create one Gem for a task you do every week — feedback, email drafting, summaries",
            "2. Try Gemini in Sheets on a real dataset — formulas, charts, =AI() classification",
            "3. Create a NotebookLM notebook for one module — upload notes, tutorials, past papers",
        ]
    )

    build_content_slide(prs,
        "Action Planning",
        [
            "What task would you like to automate or build an AI assistant for?",
            "Which tool would you use? (Gem / Apps Script / Workspace Flow / Vertex AI)",
            "What data sources would it need? (FAQ doc, grade sheet, schedule...)",
            "What's the first step you'd take?",
        ],
        footer_text="Write it down. You don't need to build it today — but writing it down makes it happen."
    )

    build_content_slide(prs,
        "Resources & Next Steps",
        [
            "Google Workspace Learning Path: cloud.google.com/training/workspace",
            "Gemini Prompting Guide: workspace.google.com/resources/gemini-for-workspace-prompting-guide",
            "NotebookLM: notebooklm.google.com",
            "Apps Script docs: developers.google.com/apps-script",
            "Remaining GCP labs: Advanced Prompt Architectures + RAG with LangChain",
            "Vertex AI docs: cloud.google.com/vertex-ai/docs",
        ]
    )

    build_segue_slide(prs, gw_title_bg,
        "Thank you!",
        "Please complete the feedback survey."
    )

    # ── Save ──────────────────────────────────────────────────
    output_path = os.path.join(BASE_DIR, "Workshop-Full-Day-Slides.pptx")
    prs.save(output_path)
    print(f"\nSaved: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
